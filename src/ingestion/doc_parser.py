# doc_parser.py
# PDF 텍스트 레이어 파싱 + PPT 텍스트 파싱


import fitz
from pptx import Presentation

def extract_pdf_text(pdf_path: str) -> list[str]:
    """PDF 텍스트 레이어를 페이지 단위로 추출하여 list[str]로 반환"""
    doc = fitz.open(pdf_path)
    page_texts = []

    try:
        for page in doc:
            txt = page.get_text("text")
            if txt:
                page_texts.append(txt.strip())
            else:
                page_texts.append("")  # 텍스트 없는 페이지는 빈 문자열
    finally:
        doc.close()

    return page_texts

def extract_ppt_text(ppt_path: str) -> list[str]:
    """PPT/PPTX의 각 슬라이드 텍스트를 list[str]로 반환"""
    prs = Presentation(ppt_path)
    slide_texts = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        slide_texts.append("\n".join(slide_text))

    return slide_texts

